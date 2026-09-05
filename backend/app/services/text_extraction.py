"""
Extract plain text from uploaded training material so it can be fed to
Gemini for quiz generation. Supports PDF, PPTX, and DOCX.
"""
import io

import pdfplumber
import docx
from pptx import Presentation


def extract_text(filename: str, file_bytes: bytes) -> str:
    lower = (filename or "").lower()
    if lower.endswith(".pdf"):
        return _extract_pdf(file_bytes)
    if lower.endswith(".pptx"):
        return _extract_pptx(file_bytes)
    if lower.endswith(".docx"):
        return _extract_docx(file_bytes)
    raise ValueError("Unsupported file type -- upload a .pdf, .pptx, or .docx file")


def _extract_pdf(file_bytes: bytes) -> str:
    text_parts = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
    return "\n".join(text_parts)


def _extract_pptx(file_bytes: bytes) -> str:
    text_parts = []
    presentation = Presentation(io.BytesIO(file_bytes))
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_parts.append(shape.text)
    return "\n".join(text_parts)


def _extract_docx(file_bytes: bytes) -> str:
    document = docx.Document(io.BytesIO(file_bytes))
    return "\n".join(p.text for p in document.paragraphs if p.text)
